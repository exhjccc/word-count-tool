import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import re
import json
from bs4 import BeautifulSoup
from docx import Document
import pdfplumber
import markdown
from openpyxl import load_workbook
import csv
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
from xml.etree.ElementTree import parse as parse_xml
import yaml
import zipfile
import speech_recognition as sr
import os
#from rpy2.robjects import r
#from rpy2.robjects.packages import importr
import pandas as pd


def update_statistics():
    """更新统计信息"""
    text = text_input.get("1.0", tk.END)

    total_characters = len(text)
    characters_no_spaces = len(text.replace(" ", "").replace("\n", ""))
    words = re.findall(r'\b\w+\b', text)
    total_words = len(words)
    sentences = re.split(r'[.!?]', text)
    total_sentences = len([s for s in sentences if s.strip()])

    char_count_label.config(text=f"字符总数（包括空格）：{total_characters}")
    char_no_space_label.config(text=f"字符总数（不包括空格）：{characters_no_spaces}")
    word_count_label.config(text=f"单词总数：{total_words}")
    sentence_count_label.config(text=f"句子总数：{total_sentences}")


def load_file():
    """加载多种格式文件内容到文本框"""
    file_path = filedialog.askopenfilename(
        title="选择文件",
        filetypes=[
            ("所有支持的文件",
             "*.txt *.docx *.pdf *.html *.json *.ipynb *.xlsx *.md *.csv *.rtf *.tex *.pptx *.epub *.xml *.yaml *.log *.zip *.mp3 *.wav"),
            ("文本文件", "*.txt"),
            ("Word 文档", "*.docx"),
            ("PDF 文件", "*.pdf"),
            ("HTML 文件", "*.html"),
            ("JSON 文件", "*.json"),
            ("Jupyter Notebook 文件", "*.ipynb"),
            ("Excel 文件", "*.xlsx"),
            ("Markdown 文件", "*.md"),
            ("CSV 文件", "*.csv"),
            ("RTF 文件", "*.rtf"),
            ("LaTeX 文件", "*.tex"),
            ("PowerPoint 文件", "*.pptx"),
            ("EPUB 文件", "*.epub"),
            ("XML 文件", "*.xml"),
            ("YAML 文件", "*.yaml *.yml"),
            ("日志文件", "*.log"),
            ("压缩文件", "*.zip"),
            ("音频文件", "*.mp3 *.wav")
        ]
    )
    if not file_path:
        return  # 用户取消选择

    content = ""
    try:
        if file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as file:
                content = file.read()
        elif file_path.endswith(".docx"):
            content = read_docx(file_path)
        elif file_path.endswith(".pdf"):
            content = read_pdf(file_path)
        elif file_path.endswith(".html"):
            content = read_html(file_path)
        elif file_path.endswith(".json"):
            content = read_json(file_path)
        elif file_path.endswith(".ipynb"):
            content = read_ipynb(file_path)
        elif file_path.endswith(".xlsx"):
            content = read_excel(file_path)
        elif file_path.endswith(".md"):
            content = read_markdown(file_path)
        elif file_path.endswith(".csv"):
            content = read_csv(file_path)
        elif file_path.endswith(".rtf"):
            content = read_rtf(file_path)
        elif file_path.endswith(".tex"):
            content = read_latex(file_path)
        elif file_path.endswith(".pptx"):
            content = read_pptx(file_path)
        elif file_path.endswith(".epub"):
            content = read_epub(file_path)
        elif file_path.endswith(".xml"):
            content = read_xml(file_path)
        elif file_path.endswith(".yaml") or file_path.endswith(".yml"):
            content = read_yaml(file_path)
        elif file_path.endswith(".log"):
            content = read_log(file_path)
        elif file_path.endswith(".zip"):
            content = read_zip(file_path)
        elif file_path.endswith(".mp3") or file_path.endswith(".wav"):
            content = read_audio(file_path)
        elif file_path.endswith(".R"):
            content = read_r_script(file_path)
        #elif file_path.endswith(".RData") or file_path.endswith(".rda"):
            #content = read_rdata(file_path)
        #elif file_path.endswith(".rds"):
            #content = read_rds(file_path)
        elif file_path.endswith(".dta"):
            content = read_stata(file_path)
        else:
            messagebox.showerror("错误", "不支持的文件格式")
            return

        text_input.delete("1.0", tk.END)
        text_input.insert("1.0", content)
        update_statistics()
    except Exception as e:
        messagebox.showerror("错误", f"无法加载文件：{e}")


# 各种文件的处理方法

def read_docx(file_path): return "\n".join([p.text for p in Document(file_path).paragraphs])


def read_pdf(file_path): return "\n".join([p.extract_text() for p in pdfplumber.open(file_path).pages])


def read_html(file_path): return BeautifulSoup(open(file_path, encoding="utf-8"), "html.parser").get_text()


def read_json(file_path): return extract_text_from_json(json.load(open(file_path, encoding="utf-8")))


def extract_text_from_json(data): return " ".join([extract_text_from_json(v) for v in data.values()]) if isinstance(
    data, dict) else (data if isinstance(data, str) else "")


def read_ipynb(file_path): return "\n".join(
    ["\n".join(cell.get("source", "")) for cell in json.load(open(file_path, encoding="utf-8")).get("cells", [])])


def read_excel(file_path): return "\n".join(
    [" ".join(map(str, row)) for sheet in load_workbook(file_path).sheetnames for row in
     load_workbook(file_path)[sheet].iter_rows(values_only=True)])


def read_markdown(file_path): return BeautifulSoup(markdown.markdown(open(file_path, encoding="utf-8").read()),
                                                   "html.parser").get_text()


def read_csv(file_path): return "\n".join([" ".join(row) for row in csv.reader(open(file_path, encoding="utf-8"))])


def read_rtf(file_path): return rtf_to_text(open(file_path, encoding="utf-8").read())


def read_latex(file_path): return re.sub(r"\\.*?{|}|%.*", "", open(file_path, encoding="utf-8").read())


def read_pptx(file_path): return "\n".join(
    [shape.text for slide in Presentation(file_path).slides for shape in slide.shapes if hasattr(shape, "text")])



def read_xml(file_path): return "\n".join([elem.text for elem in parse_xml(file_path).getroot().iter() if elem.text])


def read_yaml(file_path): return yaml.safe_load(open(file_path, encoding="utf-8"))


def read_log(file_path): return open(file_path, "r", encoding="utf-8").read()



def read_audio(file_path):
    recognizer = sr.Recognizer()
    with sr.AudioFile(file_path) as source: return recognizer.recognize_google(recognizer.record(source))


def read_zip(file_path):
    content = ""
    try:
        # 解压 ZIP 文件
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            extract_path = file_path.replace(".zip", "_extracted")
            zip_ref.extractall(extract_path)

        # 遍历解压后的文件
        for root, dirs, files in os.walk(extract_path):
            for file in files:
                full_path = os.path.join(root, file)
                # 根据文件扩展名调用适当的读取方法
                if file.endswith(".txt"):
                    with open(full_path, "r", encoding="utf-8") as f:
                        content += f.read() + "\n"
                elif file.endswith(".docx"):
                    content += read_docx(full_path) + "\n"
                elif file.endswith(".pdf"):
                    content += read_pdf(full_path) + "\n"
                elif file.endswith(".html"):
                    content += read_html(full_path) + "\n"
                elif file.endswith(".json"):
                    content += read_json(full_path) + "\n"
                elif file.endswith(".md"):
                    content += read_markdown(full_path) + "\n"
                elif file.endswith(".rtf"):
                    content += read_rtf(full_path) + "\n"
                elif file.endswith(".xml"):
                    content += read_xml(full_path) + "\n"

        # 删除解压后的临时文件（可选）
        import shutil
        shutil.rmtree(extract_path)
    except Exception as e:
        content = f"ZIP 文件处理错误：{e}"
    return content


def read_epub(file_path):
    content = ""
    try:
        # 解压 EPUB 文件
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            extract_path = file_path.replace(".epub", "_extracted")
            zip_ref.extractall(extract_path)

        # 查找 HTML/XHTML 内容文件
        for root, dirs, files in os.walk(extract_path):
            for file in files:
                if file.endswith(".html") or file.endswith(".xhtml"):
                    full_path = os.path.join(root, file)
                    with open(full_path, "r", encoding="utf-8") as html_file:
                        soup = BeautifulSoup(html_file, "html.parser")
                        content += soup.get_text()

        # 删除解压后的临时文件（可选）
        import shutil
        shutil.rmtree(extract_path)
    except Exception as e:
        content = f"EPUB 文件处理错误：{e}"
    return content


#R
def read_r_script(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()

'''
def read_rdata(file_path):
    try:
        # 加载 R 数据文件
        r.load(file_path)
        # 获取当前环境中的所有变量名
        variables = list(r.ls())
        content = ""
        for var in variables:
            obj = r[var]  # 获取变量
            content += f"{var}: {str(obj)}\n"
        return content
    except Exception as e:
        return f"RData 文件处理错误：{e}"


def read_rds(file_path):
    try:
        # 加载 RDS 文件
        obj = r.readRDS(file_path)
        return str(obj)
    except Exception as e:
        return f"RDS 文件处理错误：{e}"

'''

#Stata
def read_stata(file_path):
    try:
        # 读取 .dta 文件
        df = pd.read_stata(file_path)
        # 将数据转换为纯文本，提取列名和前几行数据
        content = "Columns:\n" + "\n".join(df.columns) + "\n\nFirst few rows:\n"
        content += df.head().to_string(index=False)
        return content
    except Exception as e:
        return f"Stata 文件处理错误：{e}"

# GUI 主程序保持不变
# 创建主窗口
root = tk.Tk()
root.title("万能字数统计工具")

# 创建菜单
menu_bar = tk.Menu(root)
file_menu = tk.Menu(menu_bar, tearoff=0)
file_menu.add_command(label="加载文件", command=load_file)
menu_bar.add_cascade(label="文件", menu=file_menu)
root.config(menu=menu_bar)

# 创建文本框
text_input = tk.Text(root, wrap=tk.WORD, height=15, width=50, font=("Arial", 12))
text_input.pack(pady=10)

# 绑定实时更新
text_input.bind("<KeyRelease>", lambda event: update_statistics())

# 创建显示统计结果的标签
char_count_label = ttk.Label(root, text="字符总数（包括空格）：0", font=("Arial", 10))
char_count_label.pack(anchor="w", padx=10)

char_no_space_label = ttk.Label(root, text="字符总数（不包括空格）：0", font=("Arial", 10))
char_no_space_label.pack(anchor="w", padx=10)

word_count_label = ttk.Label(root, text="单词总数：0", font=("Arial", 10))
word_count_label.pack(anchor="w", padx=10)

sentence_count_label = ttk.Label(root, text="句子总数：0", font=("Arial", 10))
sentence_count_label.pack(anchor="w", padx=10)

# 初始化统计
update_statistics()

# 运行主循环
root.mainloop()